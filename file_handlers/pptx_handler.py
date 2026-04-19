import shutil
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from file_handlers.base import masked_output_path, mask_texts
from models import ProcessResult


def _merge_numeric(totals: dict, values: dict) -> None:
    for k, v in values.items():
        totals[k] = totals.get(k, 0) + v


def _iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _collect_runs(tf, loc: str, tasks: list):
    """text_frame 内の全 run を tasks に積む。各エントリは (run, loc, text)。"""
    for para in tf.paragraphs:
        for run in para.runs:
            if not run.text or len(run.text.strip()) <= 1:
                continue
            tasks.append((run, loc, run.text))


def process_pptx(
    path: Path,
    model: str,
    lm_studio_url: str,
    enabled_layers: set[str] | None = None,
    excluded_tags: set[str] | None = None,
    max_workers: int = 1,
) -> ProcessResult:
    output = masked_output_path(path)
    shutil.copy2(path, output)

    prs = Presentation(output)
    total = 0
    errors: list[str] = []
    layer_totals: dict = {}
    layer_elapsed: dict = {}
    replacements_log: list[dict] = []

    # Phase 1: collect all runs to mask
    tasks: list[tuple] = []  # (run, loc, text)
    for slide_idx, slide in enumerate(prs.slides):
        for shape in _iter_shapes(slide.shapes):
            try:
                if shape.has_text_frame:
                    loc = f"スライド{slide_idx + 1}/{shape.name}"
                    _collect_runs(shape.text_frame, loc, tasks)
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tbl = shape.table
                    for r in range(len(tbl.rows)):
                        for c in range(len(tbl.columns)):
                            loc = f"スライド{slide_idx + 1}/{shape.name}/行{r + 1}/列{c + 1}"
                            _collect_runs(tbl.cell(r, c).text_frame, loc, tasks)
            except Exception as e:
                errors.append(f"Slide{slide_idx + 1}/{shape.name}: {e}")

    if not tasks:
        prs.save(output)
        return ProcessResult(
            output_path=output,
            total_replacements=0,
            errors=errors,
            layer_totals=layer_totals,
            layer_elapsed=layer_elapsed,
            replacements_log=replacements_log,
        )

    # Phase 2: mask
    texts = [t[2] for t in tasks]
    results = mask_texts(texts, model, lm_studio_url, enabled_layers, excluded_tags, max_workers)

    # Phase 3: write back & aggregate
    for (run, loc, _orig), result in zip(tasks, results):
        try:
            run.text = result.final_text
            total += len(result.replacements)
            _merge_numeric(layer_totals, result.layer_counts)
            _merge_numeric(layer_elapsed, result.layer_elapsed)
            if result.error:
                errors.append(result.error)
            for r in result.replacements:
                replacements_log.append({**r, "location": loc})
        except Exception as e:
            errors.append(f"{loc}: {e}")

    prs.save(output)
    return ProcessResult(
        output_path=output,
        total_replacements=total,
        errors=errors,
        layer_totals=layer_totals,
        layer_elapsed=layer_elapsed,
        replacements_log=replacements_log,
    )
