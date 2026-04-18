"""
手動テスト用サンプルファイル生成スクリプト

実行方法:
  docker run --rm -v "$(pwd)":/app -w /app --entrypoint "" pii-masker \
    python tests/create_fixtures.py
"""

import sys
sys.path.insert(0, ".")

from pathlib import Path
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGBColor

OUTPUT_DIR = Path("tests/fixtures")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


# ===========================================================
# XLSX — 社員名簿
# ===========================================================

def create_xlsx_employee_list():
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "社員名簿"

    header_fill = PatternFill("solid", fgColor="2E75B6")
    header_font = Font(color="FFFFFF", bold=True)

    headers = ["社員番号", "氏名", "ふりがな", "生年月日", "電話番号", "メールアドレス", "住所", "所属部署"]
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center")

    employees = [
        ["E001", "山田太郎", "やまだたろう", "1985年4月12日", "090-1234-5678",
         "yamada.taro@example-corp.co.jp", "東京都渋谷区代々木1-2-3 代々木マンション401", "営業部"],
        ["E002", "田中花子", "たなかはなこ", "1990年11月3日", "080-9876-5432",
         "tanaka.hanako@example-corp.co.jp", "神奈川県横浜市中区山手町56-7", "人事部"],
        ["E003", "佐藤次郎", "さとうじろう", "1978年2月28日", "03-5555-1234",
         "sato.jiro@example-corp.co.jp", "〒160-0022 東京都新宿区新宿3-14-1", "開発部"],
        ["E004", "鈴木三郎", "すずきさぶろう", "1995年7月19日", "070-3333-4444",
         "suzuki.saburo@example-corp.co.jp", "千葉県船橋市本町1-1-1 船橋ビル203", "経理部"],
        ["E005", "高橋美咲", "たかはしみさき", "1988年9月30日", "06-6789-0123",
         "takahashi.misaki@example-corp.co.jp", "大阪府大阪市北区梅田2-4-9 グランフロント1501", "マーケティング部"],
    ]

    for row_idx, emp in enumerate(employees, 2):
        for col_idx, val in enumerate(emp, 1):
            ws.cell(row=row_idx, column=col_idx, value=val)

    for col in ws.columns:
        max_len = max(len(str(cell.value or "")) for cell in col)
        ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 40)

    path = OUTPUT_DIR / "sample_社員名簿.xlsx"
    wb.save(path)
    print(f"  作成: {path}")


def create_xlsx_customer_data():
    wb = openpyxl.Workbook()

    # Sheet 1: 顧客基本情報
    ws1 = wb.active
    ws1.title = "顧客基本情報"

    ws1.append(["顧客ID", "会社名", "担当者名", "電話番号", "FAX", "メール", "ウェブサイト"])
    customers = [
        ["C001", "株式会社アルファシステム", "伊藤誠一", "03-1111-2222", "03-1111-2223",
         "ito.seiichi@alpha-sys.co.jp", "https://alpha-sys.co.jp"],
        ["C002", "ベータコンサルティング合同会社", "渡辺由美", "06-3333-4444", "06-3333-4445",
         "watanabe.yumi@beta-consulting.com", "https://beta-consulting.com/contact"],
        ["C003", "ガンマ商事株式会社", "中村健太", "052-5678-9012", "052-5678-9013",
         "nakamura.kenta@gamma-shoji.jp", "https://www.gamma-shoji.jp"],
    ]
    for c in customers:
        ws1.append(c)

    # Sheet 2: 契約情報
    ws2 = wb.create_sheet("契約情報")
    ws2.append(["契約番号", "顧客ID", "担当者", "連絡先", "マイナンバー（法人番号）"])
    contracts = [
        ["CON-2024-001", "C001", "伊藤誠一", "ito.seiichi@alpha-sys.co.jp", "123456789012"],
        ["CON-2024-002", "C002", "渡辺由美", "090-9999-8888", "987654321098"],
        ["CON-2024-003", "C003", "中村健太", "nakamura.kenta@gamma-shoji.jp", "111222333444"],
    ]
    for c in contracts:
        ws2.append(c)

    path = OUTPUT_DIR / "sample_顧客データ.xlsx"
    wb.save(path)
    print(f"  作成: {path}")


def create_xlsx_mixed_no_pii():
    """PIIを含まない（マスク対象なし）ファイル — 誤検知チェック用"""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "売上データ"

    ws.append(["月", "売上（万円）", "前月比（%）", "目標達成率（%）", "備考"])
    data = [
        ["2024年1月", 1250, 102.3, 95.2, "年始キャンペーン実施"],
        ["2024年2月", 1380, 110.4, 105.1, "新製品ローンチ"],
        ["2024年3月", 1420, 102.9, 108.0, "期末需要増"],
        ["2024年4月", 980, 69.0, 74.8, "新年度調整期"],
        ["2024年5月", 1100, 112.2, 83.9, "GW商戦"],
    ]
    for row in data:
        ws.append(row)

    path = OUTPUT_DIR / "sample_売上データ_PIIなし.xlsx"
    wb.save(path)
    print(f"  作成: {path}")


# ===========================================================
# PPTX — 会議資料
# ===========================================================

def create_pptx_meeting():
    prs = Presentation()

    # Slide 1: タイトル
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "2024年度 第2四半期 営業報告"
    slide1.placeholders[1].text = (
        "報告者: 営業部長 山田太郎\n"
        "連絡先: yamada.taro@example-corp.co.jp / 090-1234-5678\n"
        "日時: 2024年10月15日"
    )

    # Slide 2: 担当者一覧
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "担当者連絡先一覧"
    tf = slide2.placeholders[1].text_frame
    tf.text = "営業チーム連絡先"
    contacts = [
        "田中花子（東日本担当）: tanaka.hanako@example-corp.co.jp  080-9876-5432",
        "佐藤次郎（西日本担当）: sato.jiro@example-corp.co.jp  03-5555-1234",
        "鈴木三郎（海外担当）  : suzuki.saburo@example-corp.co.jp  070-3333-4444",
    ]
    for c in contacts:
        p = tf.add_paragraph()
        p.text = f"• {c}"

    # Slide 3: 顧客訪問リスト（テーブル）
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "主要顧客訪問記録"

    rows, cols = 5, 4
    left, top = Inches(0.5), Inches(1.5)
    width, height = Inches(9), Inches(3)
    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

    headers = ["顧客名", "担当者", "連絡先", "訪問日"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h

    visits = [
        ["株式会社アルファシステム", "伊藤誠一", "ito.seiichi@alpha-sys.co.jp", "2024/09/05"],
        ["ベータコンサルティング合同会社", "渡辺由美", "06-3333-4444", "2024/09/12"],
        ["ガンマ商事株式会社", "中村健太", "nakamura.kenta@gamma-shoji.jp", "2024/09/20"],
        ["株式会社デルタ製造", "小林正雄", "kobayashi.masao@delta-mfg.com", "2024/09/28"],
    ]
    for r, row in enumerate(visits, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    # Slide 4: 議事録
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "重要決定事項"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "本日の決定事項："
    items = [
        "新規案件: 株式会社アルファシステム様との契約（担当: 山田太郎、090-1234-5678）",
        "住所変更対応: 〒150-0001 東京都渋谷区神宮前1-2-3（田中花子）",
        "次回会議: 2024年11月1日 於 大阪府大阪市北区梅田オフィス",
        "問い合わせ先: pii-masker-support@example-corp.co.jp",
    ]
    for item in items:
        p = tf4.add_paragraph()
        p.text = f"・{item}"

    path = OUTPUT_DIR / "sample_営業報告.pptx"
    prs.save(path)
    print(f"  作成: {path}")


def create_pptx_hr():
    prs = Presentation()

    # Slide 1: タイトル
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "新入社員研修資料"
    slide1.placeholders[1].text = "人事部 高橋美咲\ntakahashi.misaki@example-corp.co.jp"

    # Slide 2: 新入社員プロフィール
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "2024年度 新入社員一覧"
    tf = slide2.placeholders[1].text_frame
    tf.text = "今年度の新入社員（要機密扱い）"
    profiles = [
        "松本陽子 / 1999年6月15日生 / 〒220-0012 神奈川県横浜市西区みなとみらい4-5-6",
        "井上裕介 / 2000年1月8日生  / 090-7777-8888 / inoue.yusuke@example-corp.co.jp",
        "木村彩香 / 1998年12月1日生  / 080-4444-5555 / 大阪府吹田市江坂町1-23-4",
    ]
    for p_text in profiles:
        p = tf.add_paragraph()
        p.text = f"▸ {p_text}"

    path = OUTPUT_DIR / "sample_新入社員研修.pptx"
    prs.save(path)
    print(f"  作成: {path}")


# ===========================================================
# DOCX — 契約書・申込書
# ===========================================================

def create_docx_application_form():
    doc = Document()

    doc.add_heading("サービス申込書", 0)
    doc.add_paragraph("以下の内容でサービスをお申し込みいたします。")
    doc.add_paragraph("")

    # 申込者情報テーブル
    doc.add_heading("申込者情報", level=1)
    table = doc.add_table(rows=10, cols=2)
    table.style = "Table Grid"

    fields = [
        ("氏名",           "山田太郎"),
        ("ふりがな",       "やまだたろう"),
        ("生年月日",       "1985年4月12日"),
        ("電話番号",       "090-1234-5678"),
        ("メールアドレス", "yamada.taro@example-corp.co.jp"),
        ("郵便番号",       "〒150-0001"),
        ("住所",           "東京都渋谷区代々木1-2-3 代々木マンション401"),
        ("会社名",         "株式会社アルファシステム"),
        ("マイナンバー",   "123456789012"),
        ("緊急連絡先",     "田中花子（妻） 080-9876-5432"),
    ]
    for i, (label, value) in enumerate(fields):
        table.cell(i, 0).text = label
        table.cell(i, 1).text = value

    doc.add_paragraph("")

    # 利用規約同意
    doc.add_heading("ご確認事項", level=1)
    doc.add_paragraph(
        "本サービスの利用規約（https://example-corp.co.jp/terms）に同意の上、"
        "上記内容にて申し込みます。"
    )
    doc.add_paragraph("ご不明な点は support@example-corp.co.jp または 03-0000-1234 までお問い合わせください。")

    doc.add_paragraph("")
    doc.add_paragraph("申込日: 2024年10月15日")
    doc.add_paragraph("署名: 山田太郎")

    path = OUTPUT_DIR / "sample_サービス申込書.docx"
    doc.save(path)
    print(f"  作成: {path}")


def create_docx_contract():
    doc = Document()

    doc.add_heading("業務委託契約書", 0)
    doc.add_paragraph("")

    doc.add_paragraph(
        "株式会社アルファシステム（以下「甲」）と佐藤次郎（以下「乙」）は、"
        "以下の条件で業務委託契約を締結する。"
    )
    doc.add_paragraph("")

    doc.add_heading("第1条（委託業務）", level=1)
    doc.add_paragraph(
        "甲は乙に対し、システム開発業務を委託する。"
        "乙の連絡先は sato.jiro@freelance.example.com または 03-5555-1234 とする。"
    )

    doc.add_heading("第2条（委託料）", level=1)
    doc.add_paragraph(
        "甲は乙に対し、月額金50万円を支払う。"
        "振込先: 〇〇銀行 渋谷支店 普通口座 1234567（名義: サトウジロウ）"
    )

    doc.add_heading("第3条（個人情報の取り扱い）", level=1)
    doc.add_paragraph(
        "乙は業務上知り得た以下の個人情報を厳重に管理する。\n"
        "・顧客氏名・住所・電話番号等の個人識別情報\n"
        "・マイナンバー等の特定個人情報（例: 987654321098）"
    )

    doc.add_paragraph("")
    doc.add_heading("甲（委託者）", level=2)
    info_table = doc.add_table(rows=3, cols=2)
    info_table.style = "Table Grid"
    for label, val in [
        ("会社名", "株式会社アルファシステム"),
        ("代表者", "伊藤誠一"),
        ("住所",   "〒100-0001 東京都千代田区丸の内1-1-1"),
    ]:
        r = info_table.add_row()
        r.cells[0].text = label
        r.cells[1].text = val
    # Remove added empty rows (add_table already created 3)
    for i, (label, val) in enumerate([
        ("会社名", "株式会社アルファシステム"),
        ("代表者", "伊藤誠一"),
        ("住所",   "〒100-0001 東京都千代田区丸の内1-1-1"),
    ]):
        info_table.cell(i, 0).text = label
        info_table.cell(i, 1).text = val

    doc.add_paragraph("")
    doc.add_heading("乙（受託者）", level=2)
    doc.add_paragraph("氏名: 佐藤次郎")
    doc.add_paragraph("住所: 東京都新宿区新宿3-14-1")
    doc.add_paragraph("電話: 03-5555-1234")
    doc.add_paragraph("メール: sato.jiro@freelance.example.com")

    doc.add_paragraph("")
    doc.add_paragraph("締結日: 2024年10月1日")

    path = OUTPUT_DIR / "sample_業務委託契約書.docx"
    doc.save(path)
    print(f"  作成: {path}")


def create_docx_meeting_minutes():
    doc = Document()

    doc.add_heading("議事録", 0)

    # ヘッダー情報テーブル
    meta = doc.add_table(rows=5, cols=2)
    meta.style = "Table Grid"
    for i, (k, v) in enumerate([
        ("日時",    "2024年10月15日（火）14:00〜16:00"),
        ("場所",    "東京都渋谷区代々木1-2-3 第3会議室"),
        ("主催",    "山田太郎（yamada.taro@example-corp.co.jp）"),
        ("出席者",  "田中花子、佐藤次郎、鈴木三郎、高橋美咲"),
        ("欠席者",  "なし"),
    ]):
        meta.cell(i, 0).text = k
        meta.cell(i, 1).text = v

    doc.add_paragraph("")
    doc.add_heading("議題・決定事項", level=1)

    topics = [
        ("議題1: 新規顧客獲得状況",
         "株式会社ベータコンサルティングの渡辺由美様（06-3333-4444）より引き合いあり。"
         "来週中に提案書を送付する（担当: 田中花子）。"),
        ("議題2: 個人情報管理体制の強化",
         "顧客データベース（含む氏名・住所・電話番号・マイナンバー）の"
         "アクセス権限を見直す。問い合わせ先: security@example-corp.co.jp"),
        ("議題3: 次回会議",
         "2024年11月5日（火）14:00〜 於 大阪府大阪市北区梅田2-4-9 グランフロントオフィス\n"
         "リモート参加用URL: https://meet.example.com/session/abc123"),
    ]
    for title, body in topics:
        doc.add_heading(title, level=2)
        doc.add_paragraph(body)
        doc.add_paragraph("")

    doc.add_heading("アクションアイテム", level=1)
    table = doc.add_table(rows=4, cols=3)
    table.style = "Table Grid"
    for j, h in enumerate(["項目", "担当者", "期限"]):
        table.cell(0, j).text = h
    actions = [
        ("提案書送付（ベータコンサルティング様）", "田中花子 tanaka.hanako@example-corp.co.jp", "10/22"),
        ("アクセス権限リスト作成",                 "佐藤次郎 03-5555-1234",                    "10/25"),
        ("大阪会場手配",                           "高橋美咲 takahashi.misaki@example-corp.co.jp", "10/20"),
    ]
    for i, (item, owner, due) in enumerate(actions, 1):
        table.cell(i, 0).text = item
        table.cell(i, 1).text = owner
        table.cell(i, 2).text = due

    path = OUTPUT_DIR / "sample_議事録.docx"
    doc.save(path)
    print(f"  作成: {path}")


def create_docx_no_pii():
    """PIIを含まない文書 — 誤検知チェック用"""
    doc = Document()

    doc.add_heading("技術仕様書: PII マスキングパイプライン", 0)

    doc.add_heading("概要", level=1)
    doc.add_paragraph(
        "本システムは4層パイプラインで個人情報を検出・マスキングする。"
        "Layer 1は正規表現、Layer 2はNER、Layer 3/4はLLMを使用する。"
    )

    doc.add_heading("処理フロー", level=1)
    for step in [
        "入力ファイル（xlsx/pptx/docx）を読み込む",
        "テキストをセル・段落単位で抽出する",
        "4層パイプラインでマスキング処理を行う",
        "マスキング済みテキストをファイルに書き戻す",
        "結果レポートを出力する",
    ]:
        doc.add_paragraph(f"• {step}")

    doc.add_heading("パフォーマンス指標", level=1)
    table = doc.add_table(rows=4, cols=3)
    table.style = "Table Grid"
    for j, h in enumerate(["指標", "目標値", "実測値"]):
        table.cell(0, j).text = h
    metrics = [
        ("処理速度（Layer 1）", "< 10ms/テキスト", "3ms"),
        ("処理速度（Layer 2）", "< 500ms/テキスト", "320ms"),
        ("誤検知率",           "< 5%",             "2.3%"),
    ]
    for i, (m, t, a) in enumerate(metrics, 1):
        table.cell(i, 0).text = m
        table.cell(i, 1).text = t
        table.cell(i, 2).text = a

    path = OUTPUT_DIR / "sample_技術仕様書_PIIなし.docx"
    doc.save(path)
    print(f"  作成: {path}")


# ===========================================================
# メイン
# ===========================================================

if __name__ == "__main__":
    print("手動テスト用サンプルファイルを生成中...\n")

    print("[XLSX]")
    create_xlsx_employee_list()
    create_xlsx_customer_data()
    create_xlsx_mixed_no_pii()

    print("\n[PPTX]")
    create_pptx_meeting()
    create_pptx_hr()

    print("\n[DOCX]")
    create_docx_application_form()
    create_docx_contract()
    create_docx_meeting_minutes()
    create_docx_no_pii()

    print(f"\n完了: {OUTPUT_DIR} に {len(list(OUTPUT_DIR.iterdir()))} ファイルを生成しました。")
    print("\n--- 生成ファイル一覧 ---")
    for f in sorted(OUTPUT_DIR.iterdir()):
        print(f"  {f.name}")
    print("\n--- 手動テスト手順 ---")
    print("  1. docker compose up でUIを起動")
    print("  2. ブラウザで http://localhost:8501 を開く")
    print("  3. フォルダに tests/fixtures を指定（コンテナ内パスに変換が必要な場合は docker-compose.yml を参照）")
    print("  4. 「マスキング開始」をクリックして結果を確認")
