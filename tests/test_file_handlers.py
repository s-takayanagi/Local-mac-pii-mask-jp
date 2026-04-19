import sys
sys.path.insert(0, ".")

import pytest
from pathlib import Path
from unittest.mock import patch, MagicMock

import openpyxl
from pptx import Presentation
from pptx.util import Inches
from docx import Document

from models import MaskResult


_MODEL = "test-model"
_URL = "http://localhost:1234/v1/chat/completions"

_MASK_RESULT_WITH_PII = MaskResult(
    final_text="[氏名]のメール: [メール]",
    replacements=[
        {"original": "山田太郎", "tag": "[氏名]", "layer": "layer1"},
        {"original": "test@example.com", "tag": "[メール]", "layer": "layer1"},
    ],
    confidence=1.0,
    layer_counts={"layer1": 2, "layer2": 0, "layer3": 0, "layer4": 0},
)

_MASK_RESULT_CLEAN = MaskResult(
    final_text="変更なし",
    replacements=[],
    confidence=1.0,
    layer_counts={"layer1": 0, "layer2": 0, "layer3": 0, "layer4": 0},
)


# ===========================================================
# XLSX
# ===========================================================

class TestXlsxHandler:
    def _create_xlsx(self, tmp_path: Path, cell_values: dict) -> Path:
        wb = openpyxl.Workbook()
        ws = wb.active
        for coord, val in cell_values.items():
            ws[coord] = val
        path = tmp_path / "test.xlsx"
        wb.save(path)
        return path

    def test_pii_cell_is_masked(self, tmp_path):
        path = self._create_xlsx(tmp_path, {"A1": "山田太郎のメール: test@example.com"})
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII):
            from file_handlers.xlsx_handler import process_xlsx
            result = process_xlsx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert result.total_replacements == 2
        assert result.output_path.exists()
        wb_out = openpyxl.load_workbook(result.output_path)
        assert wb_out.active["A1"].value == "[氏名]のメール: [メール]"

    def test_numeric_cell_skipped(self, tmp_path):
        path = self._create_xlsx(tmp_path, {"A1": 12345, "B1": "山田太郎のメール: test@example.com"})
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII) as mock_mask:
            from file_handlers.xlsx_handler import process_xlsx
            process_xlsx(path, _MODEL, _URL, enabled_layers={"layer1"})

        # 数値セル (A1) に対しては mask_text が呼ばれない
        for call in mock_mask.call_args_list:
            assert call[0][0] != 12345

    def test_short_string_cell_skipped(self, tmp_path):
        path = self._create_xlsx(tmp_path, {"A1": "A", "B1": "山田太郎のメール: test@example.com"})
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII) as mock_mask:
            from file_handlers.xlsx_handler import process_xlsx
            process_xlsx(path, _MODEL, _URL, enabled_layers={"layer1"})

        calls = [c[0][0] for c in mock_mask.call_args_list]
        assert "A" not in calls

    def test_output_path_has_masked_suffix(self, tmp_path):
        path = self._create_xlsx(tmp_path, {"A1": "テスト文字列"})
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_CLEAN):
            from file_handlers.xlsx_handler import process_xlsx
            result = process_xlsx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert result.output_path.stem.endswith("_masked")
        assert result.output_path != path

    def test_original_file_not_overwritten(self, tmp_path):
        path = self._create_xlsx(tmp_path, {"A1": "山田太郎のメール: test@example.com"})
        original_content = path.read_bytes()
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII):
            from file_handlers.xlsx_handler import process_xlsx
            process_xlsx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert path.read_bytes() == original_content

    def test_replacements_log_has_location(self, tmp_path):
        path = self._create_xlsx(tmp_path, {"A1": "山田太郎のメール: test@example.com"})
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII):
            from file_handlers.xlsx_handler import process_xlsx
            result = process_xlsx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert len(result.replacements_log) > 0
        for entry in result.replacements_log:
            assert "location" in entry
            assert "A1" in entry["location"]

    def test_empty_workbook(self, tmp_path):
        path = self._create_xlsx(tmp_path, {})
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_CLEAN) as mock_mask:
            from file_handlers.xlsx_handler import process_xlsx
            result = process_xlsx(path, _MODEL, _URL, enabled_layers={"layer1"})

        mock_mask.assert_not_called()
        assert result.total_replacements == 0

    def test_layer_totals_aggregated(self, tmp_path):
        path = self._create_xlsx(tmp_path, {
            "A1": "山田太郎のメール: test@example.com",
            "B1": "別のテキスト文字列",
        })
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII):
            from file_handlers.xlsx_handler import process_xlsx
            result = process_xlsx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert result.layer_totals.get("layer1", 0) >= 2


# ===========================================================
# PPTX
# ===========================================================

class TestPptxHandler:
    def _create_pptx(self, tmp_path: Path, texts: list[str]) -> Path:
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        for i, text in enumerate(texts):
            txBox = slide.shapes.add_textbox(Inches(i), Inches(0), Inches(3), Inches(1))
            txBox.text_frame.text = text
        path = tmp_path / "test.pptx"
        prs.save(path)
        return path

    def test_text_frame_masked(self, tmp_path):
        path = self._create_pptx(tmp_path, ["山田太郎のメール: test@example.com"])
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII):
            from file_handlers.pptx_handler import process_pptx
            result = process_pptx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert result.total_replacements == 2
        assert result.output_path.exists()

    def test_output_path_has_masked_suffix(self, tmp_path):
        path = self._create_pptx(tmp_path, ["テスト"])
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_CLEAN):
            from file_handlers.pptx_handler import process_pptx
            result = process_pptx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert result.output_path.stem.endswith("_masked")

    def test_original_not_overwritten(self, tmp_path):
        path = self._create_pptx(tmp_path, ["山田太郎"])
        original = path.read_bytes()
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII):
            from file_handlers.pptx_handler import process_pptx
            process_pptx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert path.read_bytes() == original

    def test_replacements_log_has_slide_location(self, tmp_path):
        path = self._create_pptx(tmp_path, ["山田太郎のメール: test@example.com"])
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII):
            from file_handlers.pptx_handler import process_pptx
            result = process_pptx(path, _MODEL, _URL, enabled_layers={"layer1"})

        for entry in result.replacements_log:
            assert "スライド" in entry["location"]

    def test_short_text_skipped(self, tmp_path):
        path = self._create_pptx(tmp_path, ["A"])
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII) as mock_mask:
            from file_handlers.pptx_handler import process_pptx
            process_pptx(path, _MODEL, _URL, enabled_layers={"layer1"})

        mock_mask.assert_not_called()


# ===========================================================
# DOCX
# ===========================================================

class TestDocxHandler:
    def _create_docx(self, tmp_path: Path, paragraphs: list[str]) -> Path:
        doc = Document()
        for text in paragraphs:
            doc.add_paragraph(text)
        path = tmp_path / "test.docx"
        doc.save(path)
        return path

    def _create_docx_with_table(self, tmp_path: Path, table_data: list[list[str]]) -> Path:
        doc = Document()
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 1
        table = doc.add_table(rows=rows, cols=cols)
        for r, row_data in enumerate(table_data):
            for c, cell_text in enumerate(row_data):
                table.cell(r, c).text = cell_text
        path = tmp_path / "test.docx"
        doc.save(path)
        return path

    def test_paragraph_pii_masked(self, tmp_path):
        path = self._create_docx(tmp_path, ["山田太郎のメール: test@example.com"])
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII):
            from file_handlers.docx_handler import process_docx
            result = process_docx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert result.total_replacements == 2
        assert result.output_path.exists()

    def test_output_path_has_masked_suffix(self, tmp_path):
        path = self._create_docx(tmp_path, ["テスト"])
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_CLEAN):
            from file_handlers.docx_handler import process_docx
            result = process_docx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert result.output_path.stem.endswith("_masked")

    def test_original_not_overwritten(self, tmp_path):
        path = self._create_docx(tmp_path, ["山田太郎"])
        original = path.read_bytes()
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII):
            from file_handlers.docx_handler import process_docx
            process_docx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert path.read_bytes() == original

    def test_table_cell_masked(self, tmp_path):
        path = self._create_docx_with_table(tmp_path, [
            ["氏名", "山田太郎"],
            ["メール", "test@example.com"],
        ])
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII):
            from file_handlers.docx_handler import process_docx
            result = process_docx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert result.total_replacements > 0

    def test_unmasked_name_field_generates_warning(self, tmp_path):
        # 氏名ラベルの隣のセルが未マスクの場合、警告が出る
        path = self._create_docx_with_table(tmp_path, [
            ["氏名", "山田太郎"],
        ])
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_CLEAN):
            from file_handlers.docx_handler import process_docx
            result = process_docx(path, _MODEL, _URL, enabled_layers={"layer1"})

        assert len(result.warnings) > 0
        assert any("氏名" in w for w in result.warnings)

    def test_replacements_log_has_paragraph_location(self, tmp_path):
        path = self._create_docx(tmp_path, ["山田太郎のメール: test@example.com"])
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII):
            from file_handlers.docx_handler import process_docx
            result = process_docx(path, _MODEL, _URL, enabled_layers={"layer1"})

        for entry in result.replacements_log:
            assert "段落" in entry["location"]

    def test_empty_paragraph_skipped(self, tmp_path):
        path = self._create_docx(tmp_path, [""])
        with patch("core.pipeline.mask_text", return_value=_MASK_RESULT_WITH_PII) as mock_mask:
            from file_handlers.docx_handler import process_docx
            process_docx(path, _MODEL, _URL, enabled_layers={"layer1"})

        mock_mask.assert_not_called()
